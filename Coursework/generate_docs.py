from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

# Generate Word document from markdown
report_md = 'Project_Report.md'
doc = Document()
with open(report_md, 'r', encoding='utf-8') as f:
    for line in f:
        line = line.rstrip('\n')
        if not line:
            continue
        if line.startswith('# '):
            doc.add_heading(line[2:], level=0)
        elif line.startswith('## '):
            doc.add_heading(line[3:], level=1)
        elif line.startswith('### '):
            doc.add_heading(line[4:], level=2)
        elif line.startswith('- '):
            doc.add_paragraph(line[2:], style='List Bullet')
        else:
            doc.add_paragraph(line)
doc.save('CoU-MIS_Project_Report.docx')

# Generate PowerPoint presentation from outline
outline_md = 'presentation_outline.md'
prs = Presentation()
with open(outline_md, 'r', encoding='utf-8') as f:
    title = None
    bullets = []
    for raw in f:
        line = raw.strip()
        if not line:
            continue
        if line.startswith('## '):
            if title is not None:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = title
                body = slide.shapes.placeholders[1].text_frame
                if bullets:
                    body.text = bullets[0]
                    for bullet in bullets[1:]:
                        p = body.add_paragraph()
                        p.text = bullet
                        p.level = 0
                bullets = []
            title = line[3:]
        elif line.startswith('- '):
            bullets.append(line[2:])
        else:
            bullets.append(line)
    if title is not None:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        body = slide.shapes.placeholders[1].text_frame
        if bullets:
            body.text = bullets[0]
            for bullet in bullets[1:]:
                p = body.add_paragraph()
                p.text = bullet
                p.level = 0

# Add a visual architecture slide at the end
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = 'CoU-MIS Architecture Overview'
left = Inches(1)
top = Inches(1.8)
width = Inches(3)
height = Inches(1.2)
shape1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
shape1.text = 'Users and Stakeholders\n(Bishop, Priests, Treasurer)'
left2 = Inches(4.5)
shape2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left2, top, width, height)
shape2.text = 'Database Backend\n(PostgreSQL schema)'
left3 = Inches(8)
shape3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left3, top, width, height)
shape3.text = 'Reports and Dashboard\n(Financial & Sacramental)'
prs.save('CoU-MIS_Presentation.pptx')
